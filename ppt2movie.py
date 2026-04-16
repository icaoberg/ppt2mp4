#!/usr/bin/env python3
"""
ppt2movie — PPT to Video converter with AI voice narration.

Converts a PowerPoint presentation (.pptx/.ppt) to a video file where a
Microsoft Neural TTS voice reads the presenter notes for each slide.
Slides without notes are displayed silently for a configurable duration.

Supported output formats (selected automatically from the output filename):
    .mp4  .m4v            H.264 + AAC  (default, widest compatibility)
    .ogv  .oggv           Theora + Vorbis  (open format)
    .avi                  H.264 + MP3  (legacy compatibility)
    .wmv                  WMV2 + WMA2  (Windows Media)

Python dependencies:
    pip install python-pptx edge-tts "moviepy<2" pdf2image pillow

System dependencies:
    ffmpeg:
        Ubuntu/Debian:  sudo apt-get install ffmpeg libavcodec-extra
        macOS:          brew install ffmpeg
        Windows:        https://ffmpeg.org/download.html

    LibreOffice (converts PPTX → PDF):
        Ubuntu/Debian:  sudo apt-get install libreoffice
        macOS:          brew install --cask libreoffice
        Windows:        https://www.libreoffice.org/download/

    Poppler (converts PDF → images):
        Ubuntu/Debian:  sudo apt-get install poppler-utils
        macOS:          brew install poppler
        Windows:        https://github.com/oschwartz10612/poppler-windows

Usage:
    python ppt2movie.py presentation.pptx
    python ppt2movie.py presentation.pptx --voice en-GB-RyanNeural --silent 5
    python ppt2movie.py presentation.pptx --output ~/Videos/talk.mp4
"""

import sys
import os
import asyncio
import tempfile
import subprocess
import shutil
import argparse
from datetime import date
from pathlib import Path

# ---------------------------------------------------------------------------
# Dependency checks with helpful error messages
# ---------------------------------------------------------------------------

def _require(package: str, import_name: str | None = None):
    name = import_name or package
    try:
        __import__(name)
    except ImportError:
        print(f"Missing package '{package}'. Install it with:\n  pip install {package}\n")
        sys.exit(1)

_require("pptx", "pptx")
_require("edge_tts", "edge_tts")
_require("moviepy", "moviepy")
_require("pdf2image", "pdf2image")
_require("PIL", "PIL")

from pptx import Presentation           # noqa: E402
import edge_tts                         # noqa: E402
from moviepy.editor import (            # noqa: E402
    ImageClip, AudioFileClip, concatenate_videoclips
)
from pdf2image import convert_from_path # noqa: E402
from PIL import Image                   # noqa: E402

# ---------------------------------------------------------------------------
# Defaults
# ---------------------------------------------------------------------------

DEFAULT_VOICE_MICROSOFT = "en-US-ChristopherNeural"   # Microsoft neural TTS voice
DEFAULT_VOICE_ELEVENLABS = "Rachel"                   # ElevenLabs TTS voice
DEFAULT_SILENCE = 3.0                  # seconds for slides without notes
VIDEO_FPS = 24
VIDEO_WIDTH = 1920
VIDEO_HEIGHT = 1080

# Map output file extension → ffmpeg codec settings
FORMAT_MAP = {
    ".mp4":  {"codec": "libx264",    "audio_codec": "aac",        "temp_suffix": ".temp.m4a"},
    ".m4v":  {"codec": "libx264",    "audio_codec": "aac",        "temp_suffix": ".temp.m4a"},
    ".ogv":  {"codec": "libtheora",  "audio_codec": "libvorbis",  "temp_suffix": ".temp.ogg"},
    ".oggv": {"codec": "libtheora",  "audio_codec": "libvorbis",  "temp_suffix": ".temp.ogg"},
    ".avi":  {"codec": "libx264",    "audio_codec": "mp3",        "temp_suffix": ".temp.mp3"},
    ".wmv":  {"codec": "wmv2",       "audio_codec": "wmav2",      "temp_suffix": ".temp.wma"},
}
VALID_EXTENSIONS = list(FORMAT_MAP.keys())

# ---------------------------------------------------------------------------
# Step 1 – Extract presenter notes
# ---------------------------------------------------------------------------

def extract_notes(pptx_path: str) -> list[str]:
    """Return a list of presenter-note strings, one per slide (empty string if none)."""
    prs = Presentation(pptx_path)
    notes = []
    for slide in prs.slides:
        ns = slide.notes_slide
        if ns and ns.notes_text_frame:
            text = ns.notes_text_frame.text.strip()
            notes.append(text)
        else:
            notes.append("")
    return notes


# ---------------------------------------------------------------------------
# Step 2 – Convert slides to images via LibreOffice + pdf2image
# ---------------------------------------------------------------------------

def _find_libreoffice() -> str:
    for name in ("libreoffice", "soffice"):
        path = shutil.which(name)
        if path:
            return path
    raise RuntimeError(
        "LibreOffice not found. Install it:\n"
        "  Ubuntu/Debian:  sudo apt-get install libreoffice\n"
        "  macOS:          brew install --cask libreoffice\n"
        "  Windows:        https://www.libreoffice.org/download/"
    )


def convert_pptx_to_images(pptx_path: str, output_dir: str) -> list[str]:
    """Render each slide to a PNG file and return the sorted list of paths."""
    lo = _find_libreoffice()

    print("  Converting PPTX → PDF with LibreOffice …")
    result = subprocess.run(
        [lo, "--headless", "--convert-to", "pdf", "--outdir", output_dir, pptx_path],
        capture_output=True, text=True,
    )
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice failed:\n{result.stderr}")

    pdf_path = os.path.join(output_dir, Path(pptx_path).stem + ".pdf")
    if not os.path.exists(pdf_path):
        # LibreOffice sometimes sanitises the filename
        pdfs = list(Path(output_dir).glob("*.pdf"))
        if not pdfs:
            raise RuntimeError("LibreOffice did not produce a PDF file.")
        pdf_path = str(pdfs[0])

    print("  Converting PDF → images …")
    pages = convert_from_path(pdf_path, dpi=150)

    image_paths: list[str] = []
    for i, page in enumerate(pages):
        img = page.convert("RGB").resize((VIDEO_WIDTH, VIDEO_HEIGHT), Image.LANCZOS)
        img_path = os.path.join(output_dir, f"slide_{i + 1:04d}.png")
        img.save(img_path, "PNG")
        image_paths.append(img_path)

    print(f"  {len(image_paths)} slide image(s) ready.")
    return image_paths


# ---------------------------------------------------------------------------
# Step 3 – Generate TTS audio
# ---------------------------------------------------------------------------

async def _tts_microsoft(text: str, path: str, voice: str) -> None:
    communicate = edge_tts.Communicate(text, voice)
    await communicate.save(path)


def _tts_elevenlabs(text: str, path: str, voice: str) -> None:
    try:
        from elevenlabs.client import ElevenLabs
    except ImportError:
        print("Missing package 'elevenlabs'. Install it with:\n  pip install elevenlabs\n")
        sys.exit(1)

    api_key = os.environ.get("ELEVENLABS_API_KEY")
    if not api_key:
        print("Error: ELEVENLABS_API_KEY environment variable is not set.")
        sys.exit(1)

    client = ElevenLabs(api_key=api_key)

    # Resolve voice name to voice_id; fall back to treating the value as an ID.
    voice_id = voice
    voices_response = client.voices.get_all()
    for v in voices_response.voices:
        if v.name.lower() == voice.lower() or v.voice_id == voice:
            voice_id = v.voice_id
            break

    audio_chunks = client.text_to_speech.convert(
        voice_id=voice_id,
        text=text,
        model_id="eleven_multilingual_v2",
    )

    with open(path, "wb") as f:
        for chunk in audio_chunks:
            f.write(chunk)


def _audio_duration(path: str) -> float:
    clip = AudioFileClip(path)
    dur = clip.duration
    clip.close()
    return dur


async def generate_audio_files(
    notes: list[str],
    audio_dir: str,
    voice: str,
    library: str = "microsoft",
) -> list[tuple[str | None, float]]:
    """
    For each note, generate an MP3 file and return (path, duration).
    Slides without notes get (None, silence_duration).
    Uses Microsoft edge-tts when library='microsoft', ElevenLabs when library='elevenlabs'.
    """
    results: list[tuple[str | None, float]] = []
    for i, note in enumerate(notes):
        if note:
            path = os.path.join(audio_dir, f"audio_{i + 1:04d}.mp3")
            print(f"  Slide {i + 1}: generating narration ({len(note)} chars) …")
            if library == "elevenlabs":
                _tts_elevenlabs(note, path, voice)
            else:
                await _tts_microsoft(note, path, voice)
            duration = _audio_duration(path)
            results.append((path, duration))
        else:
            results.append((None, None))   # duration filled in later
    return results


# ---------------------------------------------------------------------------
# Step 4 – Assemble the video with moviepy
# ---------------------------------------------------------------------------

def assemble_video(
    image_paths: list[str],
    audio_info: list[tuple[str | None, float | None]],
    silence: float,
    output_path: str,
    author: str = "",
    group: str = "",
    center: str = "",
    copyright: str = "",
    codec: str = "libx264",
    audio_codec: str = "aac",
    temp_suffix: str = ".temp.m4a",
) -> None:
    """Combine slide images with their audio tracks into a single MP4."""
    if len(image_paths) != len(audio_info):
        raise ValueError(
            f"Slide count mismatch: {len(image_paths)} images vs "
            f"{len(audio_info)} audio entries."
        )

    clips = []
    for i, (img_path, (audio_path, duration)) in enumerate(
        zip(image_paths, audio_info)
    ):
        actual_duration = duration if duration is not None else silence
        print(
            f"  Slide {i + 1}: {actual_duration:.1f}s "
            f"({'narrated' if audio_path else 'silent'})"
        )
        clip = ImageClip(img_path, duration=actual_duration)
        if audio_path:
            clip = clip.set_audio(AudioFileClip(audio_path))
        clips.append(clip)

    print("  Concatenating clips …")
    final = concatenate_videoclips(clips, method="compose")

    print(f"  Writing {output_path} …")
    ffmpeg_params = []
    if author:
        ffmpeg_params += ["-metadata", f"artist={author}", "-metadata", f"author={author}"]
    if group:
        ffmpeg_params += ["-metadata", f"album_artist={group}", "-metadata", f"publisher={group}"]
    if center:
        ffmpeg_params += ["-metadata", f"organization={center}"]
    if copyright:
        ffmpeg_params += ["-metadata", f"copyright={copyright}"]
    ffmpeg_params += ["-metadata", f"date={date.today().isoformat()}"]
    final.write_videofile(
        output_path,
        fps=VIDEO_FPS,
        codec=codec,
        audio_codec=audio_codec,
        temp_audiofile=str(Path(output_path).with_suffix(temp_suffix)),
        remove_temp=True,
        verbose=False,
        logger=None,
        ffmpeg_params=ffmpeg_params if ffmpeg_params else None,
    )
    final.close()
    for clip in clips:
        clip.close()


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description=(
            "Convert a PowerPoint presentation to a video with AI voice narration.\n"
            "A Microsoft Neural TTS voice reads the presenter notes for each slide.\n"
            "Slides without notes are shown silently for a configurable duration."
        ),
        epilog=(
            "examples:\n"
            "  %(prog)s talk.pptx\n"
            "  %(prog)s talk.pptx --output ~/Videos/talk.mp4\n"
            "  %(prog)s talk.pptx --output talk.wmv\n"
            "  %(prog)s talk.pptx --voice en-GB-RyanNeural --silent 5\n"
            "  %(prog)s talk.pptx --author 'Jane Smith' --group 'Research Lab'\n"
            "\n"
            "supported output formats:\n"
            "  .mp4 / .m4v   H.264 + AAC      (default, widest compatibility)\n"
            "  .ogv / .oggv  Theora + Vorbis  (open format)\n"
            "  .avi          H.264 + MP3      (legacy compatibility)\n"
            "  .wmv          WMV2 + WMA2      (Windows Media)\n"
            "\n"
            "browse available voices:\n"
            "  edge-tts --list-voices"
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument("pptx_file", help="Path to the input .pptx/.ppt file.")
    parser.add_argument(
        "--library",
        choices=["microsoft", "elevenlabs"],
        default="microsoft",
        help=(
            "TTS library to use for narration (default: %(default)s). "
            "'microsoft' uses Microsoft Neural TTS via edge-tts. "
            "'elevenlabs' uses ElevenLabs TTS (requires ELEVENLABS_API_KEY env var)."
        ),
    )
    parser.add_argument(
        "--voice",
        default=None,
        metavar="NAME",
        help=(
            "TTS voice name or ID. "
            f"Default for microsoft: {DEFAULT_VOICE_MICROSOFT}. "
            f"Default for elevenlabs: {DEFAULT_VOICE_ELEVENLABS}. "
            "For Microsoft, run `edge-tts --list-voices` to see available voices. "
            "For ElevenLabs, pass a voice name or voice_id from your ElevenLabs account."
        ),
    )
    parser.add_argument(
        "--silent",
        type=float,
        default=DEFAULT_SILENCE,
        metavar="SECONDS",
        help="Duration in seconds to show slides that have no presenter notes (default: %(default)s).",
    )
    parser.add_argument(
        "--output",
        default=None,
        metavar="PATH",
        help=(
            "Output video path. The format is determined by the file extension "
            f"({', '.join(VALID_EXTENSIONS)}). "
            "Defaults to the input filename with a .mp4 extension."
        ),
    )
    parser.add_argument(
        "--author",
        default="Ivan Cao-Berg",
        metavar="NAME",
        help="Author name to embed in the video metadata (default: %(default)s).",
    )
    parser.add_argument(
        "--group",
        default="Biomedical Applications Group",
        metavar="NAME",
        help="Group name to embed in the video metadata (default: %(default)s).",
    )
    parser.add_argument(
        "--center",
        default="Pittsburgh Supercomputing Center",
        metavar="NAME",
        help="Center name to embed in the video metadata (default: %(default)s).",
    )
    parser.add_argument(
        "--copyright",
        default=f"{date.today().year} Ivan Cao-Berg at the Pittsburgh Computing Center in Carnegie Mellon University",
        metavar="TEXT",
        help="Copyright string to embed in the video metadata (default: %(default)s).",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()

    pptx_path = Path(args.pptx_file).resolve()
    if not pptx_path.exists():
        print(f"Error: file not found: {pptx_path}")
        sys.exit(1)
    if pptx_path.suffix.lower() not in (".pptx", ".ppt"):
        print(f"Error: expected a .pptx/.ppt file, got: {pptx_path.suffix}")
        sys.exit(1)

    output_path = Path(args.output).resolve() if args.output else pptx_path.with_suffix(".mp4")
    ext = output_path.suffix.lower()
    if ext not in FORMAT_MAP:
        print(f"Error: unsupported output format '{ext}'. Valid extensions: {', '.join(VALID_EXTENSIONS)}")
        sys.exit(1)
    fmt = FORMAT_MAP[ext]

    # Resolve voice default based on chosen library
    if args.voice is None:
        voice = DEFAULT_VOICE_ELEVENLABS if args.library == "elevenlabs" else DEFAULT_VOICE_MICROSOFT
    else:
        voice = args.voice

    print("=" * 60)
    print("PPT → Video converter")
    print("=" * 60)
    print(f"  Input  : {pptx_path}")
    print(f"  Output : {output_path}")
    print(f"  Library: {args.library}")
    print(f"  Voice  : {voice}")
    print(f"  Silence: {args.silent}s per slide without notes")
    print(f"  Author : {args.author}")
    print(f"  Group  : {args.group}")
    print(f"  Center : {args.center}")
    print(f"  Copyright: {args.copyright}")
    print()

    with tempfile.TemporaryDirectory(prefix="ppt2video_") as tmpdir:

        # 1. Extract notes
        print("Step 1 – Extracting presenter notes …")
        notes = extract_notes(str(pptx_path))
        n_with = sum(1 for n in notes if n)
        print(f"  {len(notes)} slide(s) total, {n_with} with presenter notes.\n")

        # 2. Slide images
        print("Step 2 – Rendering slides …")
        image_paths = convert_pptx_to_images(str(pptx_path), tmpdir)
        print()

        # Sanity-check slide count (PDF may differ from PPTX for hidden slides, etc.)
        if len(image_paths) != len(notes):
            print(
                f"  Warning: PPTX has {len(notes)} slide(s) but PDF rendered "
                f"{len(image_paths)} page(s). Truncating to the shorter list."
            )
            count = min(len(image_paths), len(notes))
            image_paths = image_paths[:count]
            notes = notes[:count]

        # 3. TTS audio
        print("Step 3 – Generating AI narration …")
        audio_info = asyncio.run(
            generate_audio_files(notes, tmpdir, voice, args.library)
        )
        print()

        # 4. Assemble video
        print("Step 4 – Assembling video …")
        assemble_video(
            image_paths, audio_info, args.silent, str(output_path),
            args.author, args.group, args.center, args.copyright,
            **fmt,
        )

    print()
    print("Done!")
    print(f"  Video saved to: {output_path}")


if __name__ == "__main__":
    main()
